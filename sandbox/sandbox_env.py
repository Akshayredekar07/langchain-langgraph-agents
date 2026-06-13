"""
Document generation LangChain tools: PDF, DOCX, PPTX from Markdown.

Each tool runs inside an isolated subprocess sandbox.
Output is uploaded to S3 with a presigned URL; falls back to local
storage if S3 is unconfigured.

Required env vars:
    DOC_S3_BUCKET_NAME           S3 bucket name
    DOC_S3_REGION_NAME           AWS region (default: us-east-1)
    DOC_S3_BUCKET_PREFIX         Key prefix (default: AgentDocumentGeneration/)
    PRESIGNED_URL_EXPIRY_SECONDS Presigned URL TTL in seconds (default: 3600)
    DOC_LOCAL_OUTPUT_DIR         Local fallback dir (default: output/generated_docs)
    DOC_LOCAL_BASE_URL           Base URL for local fallback (default: empty)
    DOC_SANDBOX_TMP_DIR          Sandbox temp root (default: output/doc_sandbox_tmp)

Optional AWS credentials:
    AWS_ACCESS_KEY_ID
    AWS_SECRET_ACCESS_KEY

PDF uses xhtml2pdf (pure Python, no system binaries).
Fonts: Georgia (body/headings) + Courier New (code) — both native, no download needed.
Install: pip install xhtml2pdf markdown python-docx python-pptx requests
"""

from __future__ import annotations

import json
import logging
import os
import re
import shutil
import subprocess
import sys
import uuid
from pathlib import Path
from typing import Any

import boto3
from botocore.config import Config as BotoConfig
from dotenv import find_dotenv, load_dotenv
from pydantic import BaseModel, Field
from langchain_core.tools import tool

env_file = os.getenv("ENV_FILE")
load_dotenv(env_file or find_dotenv(".env", usecwd=True))

logger = logging.getLogger(__name__)

S3_BUCKET_NAME       = os.getenv("DOC_S3_BUCKET_NAME")
S3_PREFIX            = os.getenv("DOC_S3_BUCKET_PREFIX", "AgentDocumentGeneration/")
PRESIGNED_URL_EXPIRY = int(os.getenv("PRESIGNED_URL_EXPIRY_SECONDS", "3600"))
AWS_DEFAULT_REGION   = os.getenv("DOC_S3_REGION_NAME", "us-east-1")
LOCAL_DOC_OUTPUT_DIR = Path(os.getenv("DOC_LOCAL_OUTPUT_DIR", "output/generated_docs"))
LOCAL_DOC_URL_PREFIX = os.getenv("DOC_LOCAL_URL_PREFIX", "/generated_docs").rstrip("/")
LOCAL_DOC_BASE_URL   = os.getenv("DOC_LOCAL_BASE_URL", "").rstrip("/")
SANDBOX_TMP_DIR      = Path(os.getenv("DOC_SANDBOX_TMP_DIR", "output/doc_sandbox_tmp"))
FONTS_DIR            = Path(os.getenv("DOC_FONTS_DIR", "output/fonts"))
SANDBOX_TMP_DIR.mkdir(parents=True, exist_ok=True)
FONTS_DIR.mkdir(parents=True, exist_ok=True)

_MAX_OUTPUT_BYTES  = 2 * 1024 * 1024
_MAX_ERROR_DISPLAY = 4_096
_SANDBOX_TIMEOUT   = 90


# ── Sandbox ────────────────────────────────────────────────────────────────────

def _run_in_sandbox(python_code: str, output_filename: str, timeout: int = _SANDBOX_TIMEOUT) -> bytes:
    safe_name = Path(output_filename).name
    if safe_name != output_filename or not safe_name:
        raise ValueError(f"output_filename must be a plain filename, got: {output_filename!r}")

    temp_dir = (SANDBOX_TMP_DIR / f"run_{uuid.uuid4().hex}").resolve()
    temp_dir.mkdir(parents=True, exist_ok=False)

    try:
        script_path = temp_dir / "gen.py"
        output_path = temp_dir / safe_name
        script_path.write_text(python_code, encoding="utf-8")
        logger.debug("[SANDBOX] script -> %s", script_path)

        try:
            proc = subprocess.Popen(
                [sys.executable, str(script_path)],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                cwd=temp_dir,
            )
            try:
                stdout_b, stderr_b = proc.communicate(timeout=timeout)
            except subprocess.TimeoutExpired:
                proc.kill()
                proc.communicate()
                raise RuntimeError(f"[SANDBOX] Script exceeded {timeout}s timeout.")
        except RuntimeError:
            raise
        except Exception as exc:
            raise RuntimeError(f"[SANDBOX] Popen failed: {exc}") from exc

        stdout = stdout_b[:_MAX_OUTPUT_BYTES].decode("utf-8", errors="replace")
        stderr = stderr_b[:_MAX_OUTPUT_BYTES].decode("utf-8", errors="replace")

        if proc.returncode != 0:
            out_snip = stdout[-_MAX_ERROR_DISPLAY:].strip() or "(empty)"
            err_snip = stderr[-_MAX_ERROR_DISPLAY:].strip() or "(empty)"
            raise RuntimeError(
                f"[SANDBOX] Exit {proc.returncode}.\nSTDOUT: {out_snip}\nSTDERR: {err_snip}"
            )

        if not output_path.exists():
            raise RuntimeError(f"[SANDBOX] '{safe_name}' was not created.")

        file_bytes = output_path.read_bytes()
        logger.info("[SANDBOX] %d bytes <- %s", len(file_bytes), safe_name)
        return file_bytes
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


# ── PDF CSS (native fonts only — Georgia + Courier New, no download) ──────────
# NOTE on code blocks:
#   xhtml2pdf / ReportLab does NOT reliably honour white-space:pre-wrap for
#   newline preservation. The workaround used in _fix_pre_blocks() below
#   replaces every <pre><code>...</code></pre> block with a <div> whose
#   inner newlines are converted to <br/> and leading spaces to &nbsp;.
#   Therefore pre/code CSS only needs font, color, background, padding.

_PDF_CSS = """
@page {
    size: A4;
    margin: 20mm 12mm 20mm 12mm;
}

* {
    margin: 0;
    padding: 0;
}

body {
    font-family: Georgia, 'Times New Roman', serif;
    font-size: 12pt;
    line-height: 1.75;
    color: #1a1a1a;
    background: #ffffff;
}

h1 {
    font-family: Georgia, 'Times New Roman', serif;
    font-size: 20pt;
    font-weight: bold;
    color: #0d0d0d;
    border-bottom: 1.5pt solid #4A90D9;
    padding-bottom: 6pt;
    margin-bottom: 14pt;
    margin-top: 0;
    line-height: 1.2;
}
h2 {
    font-family: Georgia, 'Times New Roman', serif;
    font-size: 15pt;
    font-weight: bold;
    color: #1a1a1a;
    margin-top: 20pt;
    margin-bottom: 6pt;
    padding-bottom: 3pt;
    border-bottom: 1pt solid #e8e0d8;
}
h3 {
    font-family: Georgia, 'Times New Roman', serif;
    font-size: 12pt;
    font-weight: bold;
    color: #2c2c2c;
    margin-top: 14pt;
    margin-bottom: 4pt;
}
h4, h5, h6 {
    font-family: Georgia, 'Times New Roman', serif;
    font-size: 10.5pt;
    font-weight: bold;
    color: #555555;
    margin-top: 10pt;
    margin-bottom: 3pt;
}

p {
    margin-bottom: 8pt;
    text-align: left;
}
a {
    color: #1d4ed8;
    text-decoration: underline;
}
strong {
    font-weight: bold;
}
em {
    font-style: italic;
}

ul, ol {
    padding-left: 20pt;
    margin-bottom: 8pt;
    margin-top: 2pt;
}
li {
    margin-bottom: 3pt;
    line-height: 1.65;
}

/* Inline code */
code {
    font-family: 'Courier New', Courier, monospace;
    font-size: 10.5pt;
    background: #f0f4f8;
    color: #2a5a8c;
    padding: 1pt 4pt;
    border: 0.5pt solid #c5d8ee;
}

/* Code blocks: styled via .code-block div (see _fix_pre_blocks).
   pre/code kept minimal as fallback only. */
pre {
    font-family: 'Courier New', Courier, monospace;
    font-size: 10.5pt;
    background: #f0f4f8;
    color: #1a1a1a;
    padding: 10pt 14pt;
    margin: 8pt 0;
}
pre code {
    background: none;
    color: inherit;
    padding: 0;
    border: none;
    font-size: inherit;
}

/* Primary code block rendering — newlines expanded to <br/> by preprocessor */
.code-block {
    font-family: 'Courier New', Courier, monospace;
    font-size: 10.5pt;
    background: #f0f4f8;
    color: #1a1a1a;
    padding: 10pt 14pt;
    margin: 8pt 0;
}

/* Math / formula blocks */
.math-block {
    font-family: 'Courier New', Courier, monospace;
    font-size: 10pt;
    background: #f0f4f8;
    color: #1a1a1a;
    padding: 10pt 16pt;
    margin: 8pt 0;
    border-left: 3pt solid #4A90D9;
    line-height: 1.9;
}

blockquote {
    border-left: 3pt solid #4A90D9;
    margin: 10pt 0;
    padding: 6pt 14pt;
    background: #eef4fb;
    color: #6b6560;
    font-style: italic;
}
blockquote p {
    margin: 0;
}

table {
    width: 100%;
    border-collapse: collapse;
    margin: 10pt 0;
    font-family: Georgia, 'Times New Roman', serif;
    font-size: 9pt;
}
th {
    background: #f0f4f8;
    font-size: 8pt;
    font-weight: bold;
    letter-spacing: 0.3pt;
    text-transform: uppercase;
    padding: 6pt 9pt;
    border: 0.75pt solid #c5d8ee;
    text-align: left;
    color: #2a5a8c;
}
td {
    padding: 6pt 9pt;
    border: 0.75pt solid #c5d8ee;
    vertical-align: top;
}

hr {
    border: none;
    border-top: 1pt solid #c5d8ee;
    margin: 16pt 0;
}

img {
    max-width: 100%;
}
"""


# ── Markdown pre-processor for math formulas ──────────────────────────────────

def _preprocess_math_md(md_content: str) -> str:
    """
    Converts indented (4-space) code blocks that look like math formulas
    into fenced ```math blocks, which are then converted to
    <div class="math-block"> in HTML post-processing.
    """
    lines = md_content.split("\n")
    result = []
    i = 0
    while i < len(lines):
        line = lines[i]
        if line.startswith("    ") and not line.startswith("     "):
            block_lines = []
            while i < len(lines) and (lines[i].startswith("    ") or lines[i].strip() == ""):
                if lines[i].strip():
                    block_lines.append(lines[i][4:])
                i += 1
            if block_lines:
                result.append("```math")
                result.extend(block_lines)
                result.append("```")
                result.append("")
            continue
        result.append(line)
        i += 1
    return "\n".join(result)


def _fix_pre_blocks(html: str) -> str:
    """
    xhtml2pdf does not preserve newlines inside <pre> blocks via CSS alone.
    This replaces every <pre><code ...>...</code></pre> with a <div class="code-block">
    where:
      - newlines are converted to <br/>
      - leading spaces on each line are converted to &nbsp; chains
        so indentation is preserved in the PDF.
    Also handles bare <pre>...</pre> without inner <code>.
    """
    def convert_block(inner: str) -> str:
        # unescape HTML entities that markdown may have introduced
        inner = inner.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
        lines = inner.split("\n")
        result_lines = []
        for line in lines:
            # count leading spaces and replace with &nbsp;
            stripped = line.lstrip(" ")
            n_spaces = len(line) - len(stripped)
            result_lines.append("&nbsp;" * n_spaces + stripped)
        return "<br/>".join(result_lines)

    # <pre><code class="language-xxx">...</code></pre>
    def replace_pre_code(m: re.Match) -> str:
        inner = m.group(1)
        return f'<div class="code-block">{convert_block(inner)}</div>'

    html = re.sub(
        r'<pre><code[^>]*>(.*?)</code></pre>',
        replace_pre_code,
        html,
        flags=re.DOTALL,
    )

    # bare <pre>...</pre> (no inner code tag)
    def replace_bare_pre(m: re.Match) -> str:
        inner = m.group(1)
        return f'<div class="code-block">{convert_block(inner)}</div>'

    html = re.sub(
        r'<pre>(.*?)</pre>',
        replace_bare_pre,
        html,
        flags=re.DOTALL,
    )

    return html


def _md_to_html_with_math(md_content: str) -> str:
    """
    Convert markdown to HTML with:
      - math block support (4-space indent -> .math-block div)
      - mermaid notice box
      - code block newline fix via _fix_pre_blocks()
    """
    import markdown

    processed_md = _preprocess_math_md(md_content)

    html = markdown.markdown(
        processed_md,
        extensions=["tables", "fenced_code", "sane_lists"],
    )

    # math blocks
    html = re.sub(
        r'<pre><code class="language-math">(.*?)</code></pre>',
        lambda m: f'<div class="math-block">{m.group(1)}</div>',
        html,
        flags=re.DOTALL,
    )

    # mermaid notice
    html = re.sub(
        r'<pre><code class="language-mermaid">(.*?)</code></pre>',
        lambda m: (
            '<div class="math-block" style="color:#777;font-style:italic;">'
            '[Diagram — Mermaid requires a JS engine, showing source:]<br/>'
            + m.group(1).replace("\n", "<br/>")
            + "</div>"
        ),
        html,
        flags=re.DOTALL,
    )

    # fix newlines in all remaining pre/code blocks
    html = _fix_pre_blocks(html)

    return html


# ── Script generators ──────────────────────────────────────────────────────────

def _build_pdf_script(md_content: str, output_name: str) -> str:
    """
    xhtml2pdf-based PDF generation with native Georgia + Courier New fonts.
    No font downloads required. Code block newlines are expanded to <br/>
    by _md_to_html_with_math so xhtml2pdf renders them correctly.
    """
    css_escaped  = json.dumps(_PDF_CSS)
    full_html    = _md_to_html_with_math(md_content)
    html_escaped = json.dumps(full_html)

    return f"""
import xhtml2pdf.pisa as _pisa
import io as _io
import json as _json

_css       = _json.loads({css_escaped!r})
_html_body = _json.loads({html_escaped!r})

_full_html = f\"\"\"<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8"/>
  <style>{{_css}}</style>
</head>
<body>
{{_html_body}}
</body>
</html>\"\"\"

with open({output_name!r}, "wb") as _f:
    result = _pisa.CreatePDF(
        src=_io.StringIO(_full_html),
        dest=_f,
        encoding="utf-8",
    )

if result.err:
    raise RuntimeError(f"xhtml2pdf error: {{result.err}}")

print("PDF generated")
"""


def _build_docx_script(title: str, md_content: str, output_name: str) -> str:
    return f"""
import re as _re
from docx import Document as _Doc
from docx.shared import Pt as _Pt, RGBColor as _RGB, Cm as _Cm
from docx.oxml.ns import qn as _qn
from docx.oxml import OxmlElement as _OxmlEl

_HEADING_RE  = _re.compile(r'^(#{{1,6}})\\s+(.*)')
_BULLET_RE   = _re.compile(r'^[-*+]\\s+(.*)')
_NUMBERED_RE = _re.compile(r'^\\d+\\.\\s+(.*)')
_HR_RE       = _re.compile(r'^---+\\s*$')
_FENCE_RE    = _re.compile(r'^```')
_TABLE_RE    = _re.compile(r'^\\|.*\\|')
_TABLE_SEP   = _re.compile(r'^\\|[\\s\\-:|]+\\|')
_INLINE_RE   = _re.compile(r'(\\*\\*(.+?)\\*\\*|\\*(.+?)\\*|`(.+?)`)')

def _tokenize(text):
    lines = text.strip().split('\\n')
    tokens = []
    i = 0
    in_fence = False
    code_lines = []
    lang = ''
    while i < len(lines):
        line = lines[i]
        if _FENCE_RE.match(line):
            if not in_fence:
                in_fence = True
                lang = line[3:].strip()
                code_lines = []
            else:
                in_fence = False
                tokens.append({{'type': 'code', 'lang': lang, 'text': '\\n'.join(code_lines)}})
            i += 1
            continue
        if in_fence:
            code_lines.append(line)
            i += 1
            continue
        m = _HEADING_RE.match(line)
        if m:
            tokens.append({{'type': 'heading', 'level': len(m.group(1)), 'text': m.group(2)}})
            i += 1
            continue
        if _TABLE_RE.match(line):
            rows = []
            while i < len(lines) and _TABLE_RE.match(lines[i]):
                if not _TABLE_SEP.match(lines[i]):
                    rows.append([c.strip() for c in lines[i].strip('|').split('|')])
                i += 1
            tokens.append({{'type': 'table', 'rows': rows}})
            continue
        m = _BULLET_RE.match(line)
        if m:
            tokens.append({{'type': 'bullet', 'text': m.group(1)}})
            i += 1
            continue
        m = _NUMBERED_RE.match(line)
        if m:
            tokens.append({{'type': 'numbered', 'text': m.group(1)}})
            i += 1
            continue
        if line.startswith('>'):
            tokens.append({{'type': 'quote', 'text': line.lstrip('> ').strip()}})
            i += 1
            continue
        if _HR_RE.match(line):
            tokens.append({{'type': 'hr'}})
            i += 1
            continue
        if line.strip():
            tokens.append({{'type': 'para', 'text': line.strip()}})
        i += 1
    return tokens

def _add_inline(para, text):
    pos = 0
    for m in _INLINE_RE.finditer(text):
        if m.start() > pos:
            para.add_run(text[pos:m.start()])
        full = m.group(0)
        if full.startswith('**'):
            r = para.add_run(m.group(2)); r.bold = True
        elif full.startswith('*'):
            r = para.add_run(m.group(3)); r.italic = True
        elif full.startswith('`'):
            r = para.add_run(m.group(4))
            r.font.name = 'Courier New'
            r.font.size = _Pt(10.5)
            r.font.color.rgb = _RGB(0x2A, 0x5A, 0x8C)
        pos = m.end()
    if pos < len(text):
        para.add_run(text[pos:])

def _shading(cell, hex_color):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = _OxmlEl('w:shd')
    shd.set(_qn('w:val'), 'clear')
    shd.set(_qn('w:color'), 'auto')
    shd.set(_qn('w:fill'), hex_color)
    tcPr.append(shd)

def _border_para(p, side, val='single', sz='12', space='4', color='4A90D9'):
    pPr = p._p.get_or_add_pPr()
    pBdr = _OxmlEl('w:pBdr')
    el = _OxmlEl(f'w:{{side}}')
    el.set(_qn('w:val'), val)
    el.set(_qn('w:sz'), sz)
    el.set(_qn('w:space'), space)
    el.set(_qn('w:color'), color)
    pBdr.append(el)
    pPr.append(pBdr)

doc = _Doc()
sec = doc.sections[0]
sec.top_margin    = _Cm(2.5)
sec.bottom_margin = _Cm(2.5)
sec.left_margin   = _Cm(2.0)
sec.right_margin  = _Cm(2.0)

normal = doc.styles['Normal']
normal.font.name = 'Calibri'
normal.font.size = _Pt(12)

_title_para = doc.add_paragraph()
_tr = _title_para.add_run({title!r})
_tr.bold = True
_tr.font.name = 'Calibri'
_tr.font.size = _Pt(20)
_tr.font.color.rgb = _RGB(0x0F, 0x0E, 0x0D)
_title_para.paragraph_format.space_after = _Pt(4)
_border_para(_title_para, 'bottom', sz='8', space='4', color='4A90D9')
doc.add_paragraph()

_H_COLORS = {{
    1: _RGB(0x0F, 0x0E, 0x0D),
    2: _RGB(0x1C, 0x19, 0x17),
    3: _RGB(0x29, 0x25, 0x24),
    4: _RGB(0x57, 0x53, 0x4E),
}}
_H_SIZES = {{1: _Pt(18), 2: _Pt(14), 3: _Pt(12), 4: _Pt(10.5)}}

for _tok in _tokenize({md_content!r}):
    _tt = _tok['type']
    if _tt == 'heading':
        _lv = _tok['level']
        _p = doc.add_paragraph()
        _r = _p.add_run(_tok['text'])
        _r.bold = True
        _r.font.name = 'Calibri'
        _r.font.size = _H_SIZES.get(_lv, _Pt(12))
        _r.font.color.rgb = _H_COLORS.get(_lv, _RGB(0x1C, 0x19, 0x17))
        _p.paragraph_format.space_before = _Pt(14)
        _p.paragraph_format.space_after  = _Pt(4)
    elif _tt == 'para':
        _p = doc.add_paragraph()
        _add_inline(_p, _tok['text'])
        _p.paragraph_format.space_after = _Pt(6)
        for _rn in _p.runs:
            if not _rn.bold and not _rn.italic:
                _rn.font.size = _Pt(12)
    elif _tt == 'bullet':
        _p = doc.add_paragraph(style='List Bullet')
        _add_inline(_p, _tok['text'])
        _p.paragraph_format.space_after = _Pt(2)
        for _rn in _p.runs:
            if not _rn.bold and not _rn.italic:
                _rn.font.size = _Pt(12)
    elif _tt == 'numbered':
        _p = doc.add_paragraph(style='List Number')
        _add_inline(_p, _tok['text'])
        _p.paragraph_format.space_after = _Pt(2)
        for _rn in _p.runs:
            if not _rn.bold and not _rn.italic:
                _rn.font.size = _Pt(12)
    elif _tt == 'code':
        _p = doc.add_paragraph()
        _r = _p.add_run(_tok['text'])
        _r.font.name = 'Courier New'
        _r.font.size = _Pt(10.5)
        _r.font.color.rgb = _RGB(0x1A, 0x3A, 0x5C)
        _shd_el = _OxmlEl('w:shd')
        _shd_el.set(_qn('w:val'), 'clear')
        _shd_el.set(_qn('w:color'), 'auto')
        _shd_el.set(_qn('w:fill'), 'F0F4F8')
        _p._p.get_or_add_pPr().append(_shd_el)
        _p.paragraph_format.left_indent  = _Cm(0.5)
        _p.paragraph_format.space_before = _Pt(4)
        _p.paragraph_format.space_after  = _Pt(6)
    elif _tt == 'quote':
        _p = doc.add_paragraph()
        _add_inline(_p, _tok['text'])
        _p.paragraph_format.left_indent  = _Cm(0.8)
        _p.paragraph_format.space_before = _Pt(4)
        _p.paragraph_format.space_after  = _Pt(4)
        for _rn in _p.runs:
            _rn.italic = True
            _rn.font.color.rgb = _RGB(0x78, 0x71, 0x6C)
        _border_para(_p, 'left', sz='16', space='12', color='4A90D9')
    elif _tt == 'table':
        _rows = _tok['rows']
        if _rows:
            _cols = max(len(_rw) for _rw in _rows)
            _tbl = doc.add_table(rows=len(_rows), cols=_cols)
            _tbl.style = 'Table Grid'
            for _ri, _row in enumerate(_rows):
                for _ci, _ct in enumerate(_row):
                    _cell = _tbl.cell(_ri, _ci)
                    _cell.text = ''
                    _cp = _cell.paragraphs[0]
                    _cr = _cp.add_run(_ct)
                    _cr.font.name = 'Calibri'
                    if _ri == 0:
                        _cr.bold = True
                        _cr.font.size = _Pt(9)
                        _cr.font.color.rgb = _RGB(0x2A, 0x5A, 0x8C)
                        _shading(_cell, 'F0F4F8')
                    else:
                        _cr.font.size = _Pt(10.5)
                        if _ri % 2 == 0:
                            _shading(_cell, 'F7FAFD')
            doc.add_paragraph()
    elif _tt == 'hr':
        _p = doc.add_paragraph()
        _border_para(_p, 'bottom', val='single', sz='6', space='1', color='C5D8EE')

doc.save({output_name!r})
print("DOCX generated")
"""

def _build_pptx_script(title: str, slides: list[dict[str, Any]], output_name: str) -> str:
    slides_json = json.dumps(slides)
    return f"""
import json as _json
from pptx import Presentation as _Prs
from pptx.dml.color import RGBColor as _RGB
from pptx.util import Inches as _In, Pt as _Pt
from pptx.enum.text import PP_ALIGN as _PP_ALIGN, MSO_ANCHOR as _MSO_ANCHOR

_data = _json.loads({slides_json!r})
_prs  = _Prs()
_prs.slide_width  = _In(13.333)
_prs.slide_height = _In(7.5)

# ── Colors: clean neutral palette, no heavy theming ──────────────────────────
_BLACK   = _RGB(0x1A, 0x1A, 0x1A)
_DARK    = _RGB(0x2C, 0x2C, 0x2C)
_MID     = _RGB(0x44, 0x44, 0x44)
_MUTED   = _RGB(0x77, 0x77, 0x77)
_WHITE   = _RGB(0xFF, 0xFF, 0xFF)
_BG      = _RGB(0xFF, 0xFF, 0xFF)   # plain white slides
_CODE_BG = _RGB(0xF4, 0xF4, 0xF4)  # light grey code bg
_CODE_FG = _RGB(0x1A, 0x1A, 0x1A)  # dark code text

# ── Fonts (all pre-installed on Windows/Office, zero download) ────────────────
_F_TITLE = 'Trebuchet MS'   # headings — clean modern sans
_F_BODY  = 'Segoe UI'       # body text — crisp on screen
_F_CODE  = 'Courier New'    # monospace code

# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _BG

def _tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        _In(left), _In(top), _In(width), _In(height)
    )

def _run(tf, text, size, rgb, bold=False, italic=False,
         font=None, align=_PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = font or _F_BODY
    r.font.size  = _Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = rgb
    r.font.italic = italic

def _add_para(tf, text, size, rgb, bold=False, italic=False,
              font=None, align=_PP_ALIGN.LEFT,
              space_before=4, space_after=4):
    import re as _re
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    p.space_after  = _Pt(space_after)
    # handle inline `code` segments
    parts = _re.split(r'`(.+?)`', text)
    for i, part in enumerate(parts):
        if not part:
            continue
        r = p.add_run()
        if i % 2 == 1:
            r.text = part
            r.font.name  = _F_CODE
            r.font.size  = _Pt(size - 0.5)
            r.font.color.rgb = _CODE_FG
            r.font.bold  = False
        else:
            r.text = part
            r.font.name  = font or _F_BODY
            r.font.size  = _Pt(size)
            r.font.bold  = bold
            r.font.italic= italic
            r.font.color.rgb = rgb

# ── COVER SLIDE ───────────────────────────────────────────────────────────────
_blank = _prs.slide_layouts[6]

_cover = _prs.slides.add_slide(_blank)
_set_bg(_cover)

# Title — large, centered, bold
_box = _tb(_cover, 1.0, 2.2, 11.333, 1.8)
_tf  = _box.text_frame
_tf.vertical_anchor = _MSO_ANCHOR.MIDDLE
_tf.word_wrap = True
_run(_tf, {title!r}, size=36, rgb=_BLACK, bold=True,
     font=_F_TITLE, align=_PP_ALIGN.CENTER)

# Thin separator line (shape)
_line = _cover.shapes.add_shape(1, _In(3.5), _In(4.15), _In(6.333), _In(0.02))
_line.fill.solid()
_line.fill.fore_color.rgb = _RGB(0xCC, 0xCC, 0xCC)
_line.line.fill.background()

# Subtitle
_sub = _tb(_cover, 1.0, 4.25, 11.333, 0.6)
_tf2 = _sub.text_frame
_run(_tf2, "Generated Report", size=14, rgb=_MUTED,
     italic=True, align=_PP_ALIGN.CENTER)

# ── CONTENT SLIDES ────────────────────────────────────────────────────────────
for _idx, _sd in enumerate(_data):
    _sl = _prs.slides.add_slide(_blank)
    _set_bg(_sl)

    _slide_title = _sd.get("title", "")
    _bullets     = _sd.get("bullets", [])

    # ── Slide title ───────────────────────────────────────────────────────────
    _t = _tb(_sl, 0.5, 0.28, 12.333, 0.75)
    _tf = _t.text_frame
    _tf.vertical_anchor = _MSO_ANCHOR.MIDDLE
    _run(_tf, _slide_title, size=22, rgb=_BLACK, bold=True,
         font=_F_TITLE, align=_PP_ALIGN.LEFT)

    # Thin line under title
    _ln = _sl.shapes.add_shape(1, _In(0.5), _In(1.08), _In(12.333), _In(0.018))
    _ln.fill.solid()
    _ln.fill.fore_color.rgb = _RGB(0xCC, 0xCC, 0xCC)
    _ln.line.fill.background()

    # ── Content area ──────────────────────────────────────────────────────────
    import re as _re2
    _is_code = any(
        b.strip().startswith("```") or b.count("`") >= 3
        for b in _bullets
    )

    if _is_code:
        # Code block slide: monospaced box with light grey bg
        _cbox = _sl.shapes.add_shape(1, _In(0.5), _In(1.22), _In(12.333), _In(5.85))
        _cbox.fill.solid()
        _cbox.fill.fore_color.rgb = _CODE_BG
        _cbox.line.color.rgb = _RGB(0xDD, 0xDD, 0xDD)
        _cbox.line.width = _Pt(0.5)

        _ct = _tb(_sl, 0.7, 1.35, 11.933, 5.6)
        _ctf = _ct.text_frame
        _ctf.word_wrap = True
        _ctf.vertical_anchor = _MSO_ANCHOR.TOP

        _first = True
        for _bl in _bullets:
            _clean = _bl.strip().strip("`")
            if _first:
                p = _ctf.paragraphs[0]
                _first = False
            else:
                p = _ctf.add_paragraph()
            p.space_before = _Pt(1)
            p.space_after  = _Pt(1)
            r = p.add_run()
            r.text = _clean
            r.font.name  = _F_CODE
            r.font.size  = _Pt(11)
            r.font.color.rgb = _CODE_FG

    else:
        # Normal bullet slide
        _ct = _tb(_sl, 0.6, 1.22, 12.133, 5.9)
        _ctf = _ct.text_frame
        _ctf.word_wrap = True
        _ctf.vertical_anchor = _MSO_ANCHOR.TOP

        _first = True
        for _bl in _bullets:
            _bl = _bl.strip()
            if not _bl:
                continue

            # Sub-header: short line ending with ":"
            if _bl.endswith(":") and len(_bl) < 55:
                if _first:
                    p = _ctf.paragraphs[0]
                    _first = False
                else:
                    p = _ctf.add_paragraph()
                p.space_before = _Pt(10)
                p.space_after  = _Pt(3)
                r = p.add_run()
                r.text = _bl
                r.font.name  = _F_TITLE
                r.font.size  = _Pt(15)
                r.font.bold  = True
                r.font.color.rgb = _DARK
            else:
                if _first:
                    p = _ctf.paragraphs[0]
                    _first = False
                    p.space_before = _Pt(5)
                    p.space_after  = _Pt(5)
                    # inline code handling
                    parts = _re2.split(r'`(.+?)`', _bl)
                    _ba = False
                    for i, part in enumerate(parts):
                        if not part:
                            continue
                        r = p.add_run()
                        if i % 2 == 1:
                            r.text = part
                            r.font.name  = _F_CODE
                            r.font.size  = _Pt(13.5)
                            r.font.color.rgb = _CODE_FG
                        else:
                            r.text = ("• " + part) if not _ba else part
                            _ba = True
                            r.font.name  = _F_BODY
                            r.font.size  = _Pt(14)
                            r.font.color.rgb = _MID
                else:
                    _add_para(_ctf, "• " + _bl, size=14, rgb=_MID,
                              font=_F_BODY, space_before=5, space_after=5)

    # ── Slide number (bottom right, subtle) ───────────────────────────────────
    _pg = _tb(_sl, 12.3, 7.1, 0.85, 0.28)
    _run(_pg.text_frame, str(_idx + 2), size=9, rgb=_MUTED,
         align=_PP_ALIGN.RIGHT)

_prs.save({output_name!r})
print("PPTX generated")
"""


# ── Storage helpers ────────────────────────────────────────────────────────────

def _upload_to_s3(file_bytes: bytes, doc_type: str, profile_id: str, output_name: str) -> tuple[str, str]:
    s3_key = f"{S3_PREFIX}{doc_type}/{profile_id}/{uuid.uuid4().hex}_{output_name}"
    mime = {
        "pdf":  "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    kwargs: dict[str, Any] = {
        "region_name": AWS_DEFAULT_REGION,
        "config": BotoConfig(connect_timeout=3, read_timeout=10, retries={"max_attempts": 1}),
    }
    if os.getenv("AWS_ACCESS_KEY_ID"):
        kwargs["aws_access_key_id"]     = os.getenv("AWS_ACCESS_KEY_ID")
        kwargs["aws_secret_access_key"] = os.getenv("AWS_SECRET_ACCESS_KEY")

    s3 = boto3.client("s3", **kwargs)
    s3.put_object(
        Bucket=S3_BUCKET_NAME,
        Key=s3_key,
        Body=file_bytes,
        ContentType=mime.get(doc_type, "application/octet-stream"),
        ContentDisposition=f'attachment; filename="{output_name}"',
    )
    url = s3.generate_presigned_url(
        "get_object",
        Params={"Bucket": S3_BUCKET_NAME, "Key": s3_key},
        ExpiresIn=PRESIGNED_URL_EXPIRY,
    )
    logger.info("[S3] uploaded -> %s", s3_key)
    return s3_key, url


def _save_locally(file_bytes: bytes, doc_type: str, output_name: str) -> tuple[str, str]:
    LOCAL_DOC_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    raw_name  = Path(output_name).name or f"document.{doc_type}"
    safe_name = re.sub(r"[^A-Za-z0-9._-]+", "_", raw_name).strip("._") or f"document.{doc_type}"
    local_path = LOCAL_DOC_OUTPUT_DIR / f"{uuid.uuid4().hex}_{safe_name}"
    local_path.write_bytes(file_bytes)
    url_path = f"{LOCAL_DOC_URL_PREFIX}/{local_path.name}"
    url = f"{LOCAL_DOC_BASE_URL}{url_path}" if LOCAL_DOC_BASE_URL else url_path
    logger.info("[LOCAL] saved -> %s", local_path)
    return str(local_path), url


def _store(file_bytes: bytes, doc_type: str, profile_id: str, output_name: str) -> tuple[str, str, str]:
    if S3_BUCKET_NAME:
        try:
            key, url = _upload_to_s3(file_bytes, doc_type, profile_id, output_name)
            return "s3", key, url
        except Exception as exc:
            logger.warning("[STORE] S3 failed (%s), falling back to local.", exc)
    key, url = _save_locally(file_bytes, doc_type, output_name)
    return "local", key, url


# ── LangChain tools ────────────────────────────────────────────────────────────

class GeneratePDFInput(BaseModel):
    md_content:  str = Field(description="Full Markdown source for the document body.")
    title:       str = Field(description="Document title (rendered as H1 above the body).")
    output_name: str = Field(description="Filename with .pdf extension, e.g. 'report.pdf'.")
    profile_id:  str = Field(description="User/session ID for S3 namespacing.")


class GenerateDOCXInput(BaseModel):
    md_content:  str = Field(description="Full Markdown source for the document body.")
    title:       str = Field(description="Document title shown at the top of the DOCX.")
    output_name: str = Field(description="Filename with .docx extension, e.g. 'report.docx'.")
    profile_id:  str = Field(description="User/session ID for S3 namespacing.")


class GeneratePPTXInput(BaseModel):
    title:       str        = Field(description="Presentation title shown on the cover slide.")
    slides:      list[dict] = Field(
        description='List of slide dicts: [{"title": str, "bullets": [str, ...]}, ...]'
    )
    output_name: str        = Field(description="Filename with .pptx extension, e.g. 'deck.pptx'.")
    profile_id:  str        = Field(description="User/session ID for S3 namespacing.")


@tool(args_schema=GeneratePDFInput)
def generate_pdf(md_content: str, title: str, output_name: str, profile_id: str) -> dict:
    """
    Convert Markdown to a styled A4 PDF using xhtml2pdf (pure Python, no system binaries).
    Fonts: Georgia (body/headings) + Courier New (code) — native, no download needed.
    Math formula blocks (4-space indented) are styled distinctly from code blocks.
    Mermaid diagram blocks are noted as unsupported.

    Use when the user explicitly requests a PDF file or document download.
    Returns a dict with keys: filename, file_type, label, url, storage_key, storage.
    """
    logger.info("[TOOL] generate_pdf: %s", output_name)
    safe = Path(output_name).name
    if not safe.lower().endswith(".pdf"):
        safe += ".pdf"

    md_with_title = f"# {title}\n\n{md_content}"
    script        = _build_pdf_script(md_with_title, safe)
    fb            = _run_in_sandbox(script, safe)
    storage, key, url = _store(fb, "pdf", profile_id, safe)
    return {"filename": safe, "file_type": "PDF", "label": title, "url": url, "storage_key": key, "storage": storage}


@tool(args_schema=GenerateDOCXInput)
def generate_docx(md_content: str, title: str, output_name: str, profile_id: str) -> dict:
    """
    Convert Markdown to a styled Word DOCX.
    Use when the user explicitly requests a Word document or .docx file.
    Returns a dict with keys: filename, file_type, label, url, storage_key, storage.
    """
    logger.info("[TOOL] generate_docx: %s", output_name)
    safe = Path(output_name).name
    if not safe.lower().endswith(".docx"):
        safe += ".docx"

    script  = _build_docx_script(title, md_content, safe)
    fb      = _run_in_sandbox(script, safe)
    storage, key, url = _store(fb, "docx", profile_id, safe)
    return {"filename": safe, "file_type": "DOCX", "label": title, "url": url, "storage_key": key, "storage": storage}


@tool(args_schema=GeneratePPTXInput)
def generate_pptx(title: str, slides: list[dict], output_name: str, profile_id: str) -> dict:
    """
    Create a PowerPoint presentation from a list of slide dicts.
    Each slide dict: {"title": str, "bullets": [str, ...]}.
    Use when the user explicitly requests a slide deck or .pptx file.
    Returns a dict with keys: filename, file_type, label, url, storage_key, storage.
    """
    logger.info("[TOOL] generate_pptx: %s", output_name)
    safe = Path(output_name).name
    if not safe.lower().endswith(".pptx"):
        safe += ".pptx"

    script  = _build_pptx_script(title, slides, safe)
    fb      = _run_in_sandbox(script, safe)
    storage, key, url = _store(fb, "pptx", profile_id, safe)
    return {"filename": safe, "file_type": "PPTX", "label": title, "url": url, "storage_key": key, "storage": storage}


__all__ = [
    "generate_pdf",
    "generate_docx",
    "generate_pptx",
    "GeneratePDFInput",
    "GenerateDOCXInput",
    "GeneratePPTXInput",
]